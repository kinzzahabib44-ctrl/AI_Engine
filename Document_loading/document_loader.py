import os
import sys
import shutil
import logging
from pathlib import Path
from datetime import datetime
import re

import PyPDF2
import pandas as pd
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from docx2pdf import convert as docx2pdf_convert
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from fpdf import FPDF


DOCUMENT_DIR = "documents"
os.makedirs(DOCUMENT_DIR, exist_ok=True)
POPPLER_PATH = r"C:\Program Files\poppler-25.12.0\Library\bin"

logging.basicConfig(level=logging.INFO, format="%(message)s")

def clean_text(text: str) -> str:
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def ensure_file(src_path):
    src = Path(src_path)
    if not src.exists():
        raise FileNotFoundError(f"File not found: {src}")

    if src.is_dir():
        raise IsADirectoryError(f"Expected a file, got folder: {src}")

    dest = Path(DOCUMENT_DIR) / src.name
    if dest.exists():
        dest = dest.with_name(f"{dest.stem}_{int(datetime.now().timestamp())}{dest.suffix}")

    shutil.copy(src, dest)
    logging.info(f"Copied file to: {dest.name}")
    return dest

def ppt_to_pdf(ppt_path: Path):
    ppt = Presentation(ppt_path)
    pdf_path = ppt_path.with_suffix(".pdf")
    c = canvas.Canvas(str(pdf_path), pagesize= A4)
    width, height = A4
    margin = 40

    for slide in ppt.slides:
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        slide_text = " ".join(text)
        slide_text = clean_text(slide_text)

        lines = []
        max_chars_per_line = 90
        while slide_text:
            lines.append(slide_text[:max_chars_per_line])
            slide_text = slide_text[max_chars_per_line:]

        y = height - margin
        for line in lines:
            c.drawString(margin, y, line)
            y -= 14
            if y < margin:
                c.showPage()
                y = height - margin

        c.showPage()

    c.save()
    logging.info(f"Converted {ppt_path.name} to PDF")
    return pdf_path

def convert_to_pdf(file_path: Path):
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return file_path

    if ext == ".docx":
        sys.stdout = open(os.devnull, "w")  
        docx2pdf_convert(str(file_path), str(file_path.parent))
        sys.stdout = sys.__stdout__
        return file_path.with_suffix(".pdf")

    if ext == ".txt":
        pdf_path = file_path.with_suffix(".pdf")
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        with open(file_path, encoding="utf-8") as f:
            for line in f:
                pdf.multi_cell(0, 6, line)
        pdf.output(pdf_path)
        logging.info(f"Converted {file_path.name} to PDF")
        return pdf_path

    if ext in [".ppt", ".pptx"]:
        return ppt_to_pdf(file_path)

    raise ValueError(f"Unsupported file type: {ext}")

def iter_pdf_pages(pdf_path: Path):
    reader = PyPDF2.PdfReader(str(pdf_path))
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():
            try:
                images = convert_from_path(
                    pdf_path,
                    first_page=i,
                    last_page=i,
                    dpi=200,
                    poppler_path=POPPLER_PATH
                )
                text = pytesseract.image_to_string(images[0].convert("L"))
            except Exception as e:
                logging.warning(f"Failed to OCR page {i} of {pdf_path.name}: {e}")
                text = ""
        yield {
            "page_content": clean_text(text),
            "metadata": {
                "file_name": pdf_path.name,
                "page": i,
                "type": "pdf"
            }
        }

def load_csv(csv_path: Path):
    df = pd.read_csv(csv_path)
    docs = []
    for i, row in df.iterrows():
        text = " | ".join(f"{k}: {v}" for k, v in row.items())
        docs.append({
            "page_content": clean_text(text),
            "metadata": {
                "file_name": csv_path.name,
                "row": i + 1,
                "type": "csv"
            }
        })
    return docs
